import os
import shutil
import logging
import magic
from typing import Dict
import docx
import openpyxl
from PyPDF2 import PdfReader
import pptx
import json
import requests


class FileOrganizer:
    """文件分类器类，用于根据文件类型和内容对文件进行分类。"""

    def __init__(self, categories_file: str = "categories.json"):
        """初始化文件分类器。

        Args:
            categories_file: 分类规则配置文件路径
        """
        self.categories = self._load_categories(categories_file)
        self._setup_logging()
        self.llm_config = self.load_llm_config()

    def _load_categories(self, categories_file: str) -> Dict[str, Dict]:
        """加载分类规则。

        Args:
            categories_file: 分类规则配置文件路径

        Returns:
            分类规则字典
        """
        try:
            with open(categories_file, "r", encoding="utf-8") as f:
                return json.load(f)
        except FileNotFoundError:
            return {
                "文档": {
                    "extensions": {
                        ".txt", ".doc", ".docx", ".pdf", ".rtf", ".odt",
                        ".csv", ".tsv"
                    },
                    "keywords": {
                        "报告", "文档", "说明", "指南", "年报", "季度报告",
                        "市场概览", "分析报告"
                    }
                },
                "图片": {
                    "extensions": {
                        ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff",
                        ".ico", ".raw", ".heic"
                    },
                    "keywords": {
                        "图片", "照片", "截图", "流程图", "架构图", "设计图",
                        "UI图", "原型图"
                    }
                },
                "视频": {
                    "extensions": {
                        ".mp4", ".avi", ".mov", ".wmv", ".flv", ".mkv",
                        ".m4v", ".3gp", ".mpeg", ".mpg"
                    },
                    "keywords": {
                        "视频", "录像", "会议记录", "教程视频", "产品演示"
                    }
                },
                "音频": {
                    "extensions": {
                        ".mp3", ".wav", ".ogg", ".flac", ".aac", ".wma",
                        ".aiff", ".mid", ".midi"
                    },
                    "keywords": {
                        "音频", "音乐", "会议录音", "培训录音", "语音笔记"
                    }
                },
                "表格": {
                    "extensions": {
                        ".xls", ".xlsx", ".csv", ".ods", ".numbers"
                    },
                    "keywords": {
                        "表格", "数据表", "预算", "财务报表", "统计表"
                    }
                },
                "演示文稿": {
                    "extensions": {
                        ".ppt", ".pptx", ".odp", ".pps", ".ppsx"
                    },
                    "keywords": {
                        "演示", "幻灯片", "产品介绍", "项目汇报", "技术分享",
                        "经验分享"
                    }
                },
                "代码": {
                    "extensions": {
                        ".py", ".java", ".cpp", ".c", ".h", ".js", ".html",
                        ".css", ".php", ".rb", ".go", ".rs", ".swift",
                        ".kt", ".ts"
                    },
                    "keywords": {
                        "代码", "脚本", "程序", "配置", "源代码"
                    }
                },
                "数据": {
                    "extensions": {
                        ".json", ".xml", ".yaml", ".yml", ".sql", ".db",
                        ".sqlite", ".h5", ".parquet", ".avro", ".orc"
                    },
                    "keywords": {
                        "数据", "数据库", "数据集", "数据模型"
                    }
                },
                "设计": {
                    "extensions": {
                        ".psd", ".ai", ".sketch", ".fig", ".xd", ".indd",
                        ".eps", ".svg"
                    },
                    "keywords": {
                        "设计", "UI", "UX", "视觉设计", "平面设计"
                    }
                },
                "压缩包": {
                    "extensions": {
                        ".zip", ".rar", ".7z", ".tar", ".gz", ".bz2",
                        ".xz"
                    },
                    "keywords": {
                        "压缩", "打包", "备份"
                    }
                },
                "字体": {
                    "extensions": {
                        ".ttf", ".otf", ".woff", ".woff2", ".eot"
                    },
                    "keywords": {
                        "字体", "字库"
                    }
                }
            }

    def _setup_logging(self):
        """设置日志配置。"""
        logging.basicConfig(
            filename="file_organizer.log",
            level=logging.INFO,
            format="%(asctime)s - %(levelname)s - %(message)s"
        )

    def _extract_text_from_docx(self, file_path: str) -> str:
        """从Word文档中提取文本。

        Args:
            file_path: Word文档路径

        Returns:
            提取的文本内容
        """
        try:
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            logging.error(f"提取Word文档内容失败: {str(e)}")
            return ""

    def _extract_text_from_xlsx(self, file_path: str) -> str:
        """从Excel文件中提取文本。

        Args:
            file_path: Excel文件路径

        Returns:
            提取的文本内容
        """
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True)
            text = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.rows:
                    text.extend([str(cell.value) for cell in row if cell.value])
            return "\n".join(text)
        except Exception as e:
            logging.error(f"提取Excel文件内容失败: {str(e)}")
            return ""

    def _extract_text_from_pdf(self, file_path: str) -> str:
        """从PDF文件中提取文本。

        Args:
            file_path: PDF文件路径

        Returns:
            提取的文本内容
        """
        try:
            with open(file_path, "rb") as f:
                pdf = PdfReader(f)
                text = []
                for page in pdf.pages:
                    text.append(page.extract_text())
                return "\n".join(text)
        except Exception as e:
            logging.error(f"提取PDF文件内容失败: {str(e)}")
            return ""

    def _extract_text_from_pptx(self, file_path: str) -> str:
        """从PPT文件中提取文本。

        Args:
            file_path: PPT文件路径

        Returns:
            提取的文本内容
        """
        try:
            ppt = pptx.Presentation(file_path)
            text = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            logging.error(f"提取PPT文件内容失败: {str(e)}")
            return ""

    def _get_file_content(self, file_path: str) -> str:
        """获取文件内容。

        Args:
            file_path: 文件路径

        Returns:
            文件内容
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext in {".doc", ".docx"}:
            return self._extract_text_from_docx(file_path)
        elif file_ext in {".xls", ".xlsx"}:
            return self._extract_text_from_xlsx(file_path)
        elif file_ext == ".pdf":
            return self._extract_text_from_pdf(file_path)
        elif file_ext in {".ppt", ".pptx"}:
            return self._extract_text_from_pptx(file_path)
        return ""

    def load_llm_config(self, config_path="llm_config.json"):
        if not os.path.exists(config_path):
            return {}
        with open(config_path, "r", encoding="utf-8") as f:
            return json.load(f)

    def classify_with_llm(self, text, config):
        provider = config.get("llm_provider")
        api_key = config.get("llm_api_key")
        api_base = config.get("llm_api_base", "")
        model = config.get("llm_model", "")
        if not provider or not api_key:
            return None
        prompt = f"请判断下面内容类别（文档/表格/图片/视频/代码/合同/简历/报告/发票/课件/其他）：\n{text[:500]}"
        if provider == "openai":
            import openai
            openai.api_key = api_key
            if api_base:
                openai.api_base = api_base
            model = model or "gpt-3.5-turbo"
            resp = openai.ChatCompletion.create(
                model=model,
                messages=[{"role": "user", "content": prompt}]
            )
            return resp['choices'][0]['message']['content'].strip()
        elif provider == "deepseek":
            # DeepSeek API: https://platform.deepseek.com/
            url = api_base or "https://api.deepseek.com/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            data = {
                "model": model or "deepseek-chat",
                "messages": [{"role": "user", "content": prompt}]
            }
            resp = requests.post(url, headers=headers, json=data, timeout=30)
            return resp.json().get("choices", [{}])[0].get("message", {}).get("content", "").strip()
        elif provider == "tongyi":
            # 通义千问API: https://help.aliyun.com/zh/dashscope/developer-reference/api-details
            url = api_base or "https://dashscope.aliyuncs.com/api/v1/services/aigc/text-generation/generation"
            headers = {
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json"
            }
            data = {
                "model": model or "qwen-turbo",
                "input": {"messages": [{"role": "user", "content": prompt}]}
            }
            resp = requests.post(url, headers=headers, json=data, timeout=30)
            # 通义返回结构
            return resp.json().get("output", {}).get("choices", [{}])[0].get("message", {}).get("content", "").strip()
        return None

    def _classify_file(self, file_path: str) -> str:
        """根据文件类型和内容对文件进行分类。

        Args:
            file_path: 文件路径

        Returns:
            分类结果
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        mime = magic.Magic(mime=True)
        mime.from_file(file_path)
        # 先用扩展名
        for category, rules in self.categories.items():
            if file_ext in rules["extensions"]:
                return category
        # 用大模型API
        content = self._get_file_content(file_path)
        if content:
            llm_result = self.classify_with_llm(content, self.llm_config)
            if llm_result:
                return llm_result
            for category, rules in self.categories.items():
                if any(keyword in content for keyword in rules["keywords"]):
                    return category
        return "其他"

    def organize_files(self, directory: str):
        """组织目录中的文件。

        Args:
            directory: 要组织的目录路径
        """
        try:
            for root, _, files in os.walk(directory):
                for file in files:
                    file_path = os.path.join(root, file)
                    if os.path.isfile(file_path):
                        category = self._classify_file(file_path)
                        target_dir = os.path.join(directory, category)
                        os.makedirs(target_dir, exist_ok=True)
                        shutil.move(file_path, os.path.join(target_dir, file))
                        logging.info(
                            f"已将文件 {file} 移动到 {category} 目录"
                        )
        except Exception as e:
            logging.error(f"组织文件时发生错误: {str(e)}")


if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("使用方法: python file_organizer.py <目录路径>")
        sys.exit(1)
    organizer = FileOrganizer()
    organizer.organize_files(sys.argv[1]) 